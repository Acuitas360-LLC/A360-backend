from pptx import Presentation

ppt_path = "input.pptx"
prs = Presentation(ppt_path)

for slide_idx, slide in enumerate(prs.slides):

    print(f"\n--- Slide {slide_idx} ---")

    for shape_idx, shape in enumerate(slide.shapes):

        if shape.has_table:
            table = shape.table

            print(f"\nTable found in shape {shape_idx}")

            for r in range(len(table.rows)):
                row_values = []

                for c in range(len(table.columns)):
                    cell_text = table.cell(r, c).text.strip()
                    row_values.append(cell_text)

                print(row_values)

        if shape.has_chart:

            chart = shape.chart
            print(f"\nChart found in shape {shape_idx}")

            # Get chart title
            if chart.has_title:
                title = chart.chart_title.text_frame.text
                print("Chart Title:", title)

            # Get categories
            categories = [c.label for c in chart.plots[0].categories]
            print("Categories:", categories)

            # Get values
            for series in chart.series:
                values = list(series.values)
                print("Values:", values)